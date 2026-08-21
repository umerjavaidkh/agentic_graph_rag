"""Office documents as DocumentIR — Word, PowerPoint, Excel.

The pipeline accepted `.pdf` and nothing else, which is the wrong shape for a
real corpus: a company drive is mostly Word and PowerPoint, and a PDF export
of either loses the structure this module gets for free.

That structure is the reason these are worth having as first-class parsers
rather than converting to PDF first. A PDF forces every backend to INFER
headings from font size and boldness -- the heuristic that currently produces
"Preamble", "Smith Street" and "Dry" as section titles on a 10-K. Word states
its headings outright in the paragraph style, PowerPoint has a title
placeholder per slide, and Excel has sheet names. Nothing is guessed.

Each library is imported inside the parser that needs it, so a deployment
handling only PDFs does not carry three Office readers it never calls.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Iterator

from ...models import DKGEdge, DKGNode
from ..ir import Block, DocumentIR, PageBlock

# A heading is announced rather than inferred, so Axis-1 is told directly.
_HEADING = {"heading_hint": "heading"}


def _build(ir: DocumentIR) -> tuple[list[DKGNode], list[DKGEdge]]:
    """Shared phase 2. Imported here so the construction service is not a
    module-level dependency of a parser that may only ever be asked for IR."""
    from ...graph.construction_service import GraphConstructionService

    nodes, edges, _chunks = GraphConstructionService().build_structure(ir)
    return nodes, edges


def _page(number: int, blocks: list[Block]) -> PageBlock:
    text = "\n".join(b.text for b in blocks if b.text.strip())
    return PageBlock(
        page=number,
        text=text,
        blocks=blocks,
        regions=[b for b in blocks if b.kind in ("table", "figure")],
        confidence=1.0,        # extracted from markup, not recovered from a render
    )


def _toc_from(blocks_by_page: list[tuple[int, list[Block]]]) -> list[tuple[int, str, int]] | None:
    """The document's own outline, from headings it declared.

    Axis-1 prefers a real outline over its heading heuristics, and these
    formats can supply one truthfully. Gated at five entries with at least one
    top level, matching what the PDF path treats as a usable outline -- fewer
    than that is not a structure worth trusting.
    """
    toc = [
        (int(b.extra.get("heading_level", 1)), b.text.strip(), page)
        for page, blocks in blocks_by_page
        for b in blocks
        if b.extra.get("heading_hint") == "heading" and b.text.strip()
    ]
    if len(toc) < 5 or not any(level <= 1 for level, _, _ in toc):
        return None
    return toc


_W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


class _Para:
    __slots__ = ("paragraph",)

    def __init__(self, paragraph): self.paragraph = paragraph


class _Tbl:
    __slots__ = ("table",)

    def __init__(self, table): self.table = table


def _body_items(document) -> "list[_Para | _Tbl]":
    """Paragraphs and tables in the order they appear in the document.

    `document.paragraphs` and `document.tables` are two separate sequences,
    so reading them one after the other puts every table at the end of the
    last page no matter where it sits in the text -- wrong for reading order
    and wrong for any citation that lands on a table. The body's own child
    order is the only thing that knows where they interleave.
    """
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    items: list[_Para | _Tbl] = []
    for child in document.element.body.iterchildren():
        if child.tag == f"{_W}p":
            items.append(_Para(Paragraph(child, document)))
        elif child.tag == f"{_W}tbl":
            items.append(_Tbl(Table(child, document)))
    return items


def _paragraph_segments(para) -> list[str]:
    """One paragraph's text, split at every page break, in document order.

    Walks the XML rather than reading `paragraph.text`, because neither page
    break survives into it -- python-docx renders both as an empty string, so
    the obvious `"\f" in para.text` test never fires on a real Word file and
    every document collapses onto page 1.

    Two markers count. `w:br w:type="page"` is a break the author inserted.
    `w:lastRenderedPageBreak` is where Word itself paginated when it last
    saved: not a guess, but a record of real pagination, and the closest a
    .docx comes to having page numbers at all.

    Splitting mid-paragraph matters because that is where Word's own breaks
    usually land -- a paragraph spanning two pages should cite the page each
    half is actually on. Descending the tree also picks up text inside
    hyperlinks, which `paragraph.runs` skips.
    """
    segments = [""]
    for node in para._element.iter():
        tag = node.tag
        if tag == f"{_W}t":
            segments[-1] += node.text or ""
        elif tag == f"{_W}tab":
            segments[-1] += "\t"
        elif tag == f"{_W}lastRenderedPageBreak":
            segments.append("")
        elif tag == f"{_W}br":
            if node.get(f"{_W}type") == "page":
                segments.append("")
            else:
                segments[-1] += "\n"
        elif tag == f"{_W}cr":
            segments[-1] += "\n"
    return segments


class DocxParser:
    """Word documents.

    Word stores no page numbers -- pagination is the renderer's decision --
    so pages here are recovered from the two break markers a .docx does
    carry, and a document with neither becomes a single page, which is
    accurate rather than a guess at where pages would fall.
    """

    def parse_ir(self, source: str | Path) -> DocumentIR:
        import docx  # python-docx

        path = Path(source)
        document = docx.Document(str(path))
        pages: list[PageBlock] = []
        current: list[Block] = []
        page_no = 1

        def flush() -> None:
            """End the current page, if it has anything on it.

            A break with nothing accumulated -- one at the very top, or two
            in a row -- advances nothing, so `pages` stays contiguous with
            the page numbers on its blocks. The cost is that a deliberately
            blank page is not counted; the alternative is a `pages` list
            whose indices no longer line up with what the blocks cite.
            """
            nonlocal page_no, current
            if current:
                pages.append(_page(page_no, current))
                page_no, current = page_no + 1, []

        for item in _body_items(document):
            if isinstance(item, _Tbl):
                rows = ["\t".join(c.text.strip() for c in r.cells) for r in item.table.rows]
                body = "\n".join(r for r in rows if r.strip())
                if body:
                    current.append(
                        Block(text=body, page=page_no, source="docx", kind="table")
                    )
                continue

            para = item.paragraph
            style = (getattr(para.style, "name", "") or "")
            extra: dict[str, Any] = {}
            if style.lower().startswith("heading"):
                extra = dict(_HEADING)
                tail = style.split()[-1]
                extra["heading_level"] = int(tail) if tail.isdigit() else 1
            elif style.lower() in ("title", "subtitle"):
                extra = dict(_HEADING, heading_level=1)

            for index, segment in enumerate(_paragraph_segments(para)):
                if index:                      # a page break preceded this run
                    flush()
                text = segment.strip()
                if text:
                    current.append(
                        Block(text=text, page=page_no, source="docx", extra=dict(extra))
                    )

        pages.append(_page(page_no, current))
        ir = DocumentIR(
            source_name=path.stem,
            page_count=len(pages),
            pages=pages,
            toc=_toc_from([(p.page, p.blocks) for p in pages]),
        )
        return ir.finalize()

    def parse(self, source: str | Path) -> tuple[list[DKGNode], list[DKGEdge]]:
        return _build(self.parse_ir(source))


class PptxParser:
    """Presentations. A slide IS a page, so the mapping needs no invention,
    and the title placeholder names the section without inference."""

    def parse_ir(self, source: str | Path) -> DocumentIR:
        from pptx import Presentation

        path = Path(source)
        deck = Presentation(str(path))
        pages: list[PageBlock] = []

        for index, slide in enumerate(deck.slides, start=1):
            blocks: list[Block] = []
            title_shape = getattr(slide.shapes, "title", None)
            title_text = (getattr(title_shape, "text", "") or "").strip()
            if title_text:
                blocks.append(Block(text=title_text, page=index, source="pptx",
                                    extra=dict(_HEADING, heading_level=1)))
            for shape in slide.shapes:
                if shape is title_shape:
                    continue
                if getattr(shape, "has_table", False):
                    rows = ["\t".join(c.text.strip() for c in r.cells) for r in shape.table.rows]
                    body = "\n".join(r for r in rows if r.strip())
                    if body:
                        blocks.append(Block(text=body, page=index, source="pptx", kind="table"))
                    continue
                text = (getattr(shape, "text", "") or "").strip()
                if text:
                    blocks.append(Block(text=text, page=index, source="pptx"))
            # Speaker notes are content the slide does not show, and are often
            # where the actual argument lives.
            notes = getattr(slide, "notes_slide", None) if slide.has_notes_slide else None
            note_text = (getattr(getattr(notes, "notes_text_frame", None), "text", "") or "").strip()
            if note_text:
                blocks.append(Block(text=note_text, page=index, source="pptx",
                                    extra={"speaker_notes": True}))
            pages.append(_page(index, blocks))

        ir = DocumentIR(
            source_name=path.stem,
            page_count=len(pages),
            pages=pages,
            toc=_toc_from([(p.page, p.blocks) for p in pages]),
        )
        return ir.finalize()

    def parse(self, source: str | Path) -> tuple[list[DKGNode], list[DKGEdge]]:
        return _build(self.parse_ir(source))


class XlsxParser:
    """Workbooks as documents -- one page per sheet, the sheet name its heading.

    This is the DOCUMENT reader, not the tabular loader. A workbook whose
    sheets are really tables belongs in structured/ingestion/tabular.py, which
    turns them into a property graph with keys and relationships. This path is
    for a workbook that is prose in a grid: a checklist, a policy matrix, a
    register -- text to search, not rows to join.
    """

    def parse_ir(self, source: str | Path) -> DocumentIR:
        from openpyxl import load_workbook

        path = Path(source)
        # read_only streams rather than building the whole workbook in memory;
        # data_only takes cached formula RESULTS, since a formula string is not
        # what a reader is searching for.
        book = load_workbook(str(path), read_only=True, data_only=True)
        pages: list[PageBlock] = []

        for index, sheet in enumerate(book.worksheets, start=1):
            blocks = [Block(text=sheet.title, page=index, source="xlsx",
                            extra=dict(_HEADING, heading_level=1))]
            rows: list[str] = []
            for values in sheet.iter_rows(values_only=True):
                line = "\t".join("" if v is None else str(v) for v in values).rstrip("\t")
                if line.strip():
                    rows.append(line)
            if rows:
                blocks.append(Block(text="\n".join(rows), page=index,
                                    source="xlsx", kind="table"))
            pages.append(_page(index, blocks))

        book.close()
        ir = DocumentIR(
            source_name=path.stem,
            page_count=len(pages),
            pages=pages,
            toc=_toc_from([(p.page, p.blocks) for p in pages]),
        )
        return ir.finalize()

    def parse(self, source: str | Path) -> tuple[list[DKGNode], list[DKGEdge]]:
        return _build(self.parse_ir(source))
