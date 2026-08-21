"""Word, PowerPoint and Excel as document sources.

The pipeline accepted .pdf and nothing else. A real corpus is mostly Office
files, and a PDF export of one throws away the structure these formats state
outright -- which is the same structure the PDF path has to guess at, and
currently gets wrong ("Preamble", "Smith Street", "Dry" as section titles).
"""
import pytest

from src.unstructured.document.parser_registry import get_parser, supported_extensions

pytest.importorskip("docx")
pytest.importorskip("pptx")
pytest.importorskip("openpyxl")


def _headings(ir):
    return [b.text for p in ir.pages for b in p.blocks if b.extra.get("heading_hint") == "heading"]


@pytest.fixture
def docx_path(tmp_path):
    import docx

    d = docx.Document()
    d.add_heading("Acme Security Policy", 0)
    for title in ("Access Control", "Incident Response", "Data Retention",
                  "Vendor Review", "Training"):
        d.add_heading(title, 1)
        d.add_paragraph(f"Body text under {title}.")
    table = d.add_table(rows=1, cols=2)
    table.cell(0, 0).text, table.cell(0, 1).text = "Region", "Owner"
    path = tmp_path / "policy.docx"
    d.save(path)
    return path


def test_word_headings_are_read_not_inferred(docx_path):
    """Word states its headings in the paragraph style. Nothing is guessed
    from font size, which is what produces junk titles on the PDF path."""
    ir = get_parser(docx_path).parse_ir(docx_path)
    assert _headings(ir) == [
        "Acme Security Policy", "Access Control", "Incident Response",
        "Data Retention", "Vendor Review", "Training",
    ]


def test_word_supplies_a_real_outline(docx_path):
    """Six declared headings clear the same five-entry bar the PDF path uses
    for a usable embedded outline, so Axis-1 builds from structure."""
    ir = get_parser(docx_path).parse_ir(docx_path)
    assert ir.toc and len(ir.toc) == 6


def test_word_tables_are_regions(docx_path):
    ir = get_parser(docx_path).parse_ir(docx_path)
    assert any(b.kind == "table" for p in ir.pages for b in p.regions)


def test_a_slide_is_a_page_and_its_title_is_the_heading(tmp_path):
    """No invention needed: PowerPoint already has one title per slide."""
    from pptx import Presentation

    deck = Presentation()
    for i in (1, 2, 3):
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i} Title"
        slide.placeholders[1].text = f"Content {i}"
    path = tmp_path / "deck.pptx"
    deck.save(path)

    ir = get_parser(path).parse_ir(path)
    assert ir.page_count == 3
    assert _headings(ir) == ["Slide 1 Title", "Slide 2 Title", "Slide 3 Title"]


def test_a_sheet_is_a_page_and_its_name_is_the_heading(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    wb.active.title = "Checklist"
    wb.active.append(["Control", "Status"])
    wb.create_sheet("Register").append(["Risk", "Level"])
    path = tmp_path / "book.xlsx"
    wb.save(path)

    ir = get_parser(path).parse_ir(path)
    assert ir.page_count == 2
    assert _headings(ir) == ["Checklist", "Register"]


def test_too_few_headings_is_not_an_outline(tmp_path):
    """Two headings are not a structure worth trusting -- same bar the PDF
    path applies before preferring an outline over its heuristics."""
    from pptx import Presentation

    deck = Presentation()
    for i in (1, 2):
        deck.slides.add_slide(deck.slide_layouts[1]).shapes.title.text = f"S{i}"
    path = tmp_path / "small.pptx"
    deck.save(path)
    assert get_parser(path).parse_ir(path).toc is None


def test_the_registry_now_dispatches_office_formats():
    assert {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm"} <= supported_extensions()


def test_a_stated_heading_is_trusted_whatever_produced_it():
    """The hint used to be honoured only when block.source == 'rtldoc', so a
    parser reading headings from markup had its answer thrown away in favour
    of guessing from font size."""
    import inspect

    from src.unstructured.graph import axis1_structural

    src = inspect.getsource(axis1_structural)
    assert 'block.source == "rtldoc" and block.extra.get("heading_hint")' not in src
    assert 'block.extra.get("heading_hint") == "heading"' in src


def _docx_with_real_page_breaks(path):
    """A Word file broken the way Word breaks one -- not with a "\f".

    `add_page_break` writes `<w:br w:type="page"/>`, which is what a real
    document carries. python-docx renders it as an empty string in
    `paragraph.text`, so a parser testing for "\f" sees a single page.
    """
    from docx import Document

    document = Document()
    document.add_heading("Revenue Review", 1)
    document.add_paragraph("Net revenue was 4.2 billion in the period.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Upstream"
    table.cell(0, 1).text = "2900"
    document.add_page_break()
    document.add_heading("Risk Factors", 1)
    document.add_paragraph("Commodity price volatility remains the principal risk.")
    document.add_page_break()
    document.add_heading("Outlook", 1)
    document.add_paragraph("Guidance unchanged for the year.")
    document.save(str(path))
    return path


def test_a_real_word_page_break_starts_a_new_page(tmp_path):
    """The whole point of a page number in a citation is that it differs.

    The regression this guards is silent: every block still parses, every
    citation still renders, and every one of them says page 1.
    """
    from src.unstructured.document.office.parser import DocxParser

    ir = DocxParser().parse_ir(_docx_with_real_page_breaks(tmp_path / "breaks.docx"))

    assert ir.page_count == 3, f"page breaks not detected: {ir.page_count} page(s)"
    assert [p.page for p in ir.pages] == [1, 2, 3]
    assert "Risk Factors" in ir.pages[1].text
    assert "Outlook" in ir.pages[2].text
    assert all(b.page == p.page for p in ir.pages for b in p.blocks)


def test_a_table_is_cited_on_the_page_it_sits_on(tmp_path):
    """`document.tables` is a separate sequence from `document.paragraphs`.

    Reading one after the other puts every table on the last page, however
    early in the text it appears -- so a citation landing on a table sends
    the reader to the wrong page, and reading order is wrong too.
    """
    from src.unstructured.document.office.parser import DocxParser

    ir = DocxParser().parse_ir(_docx_with_real_page_breaks(tmp_path / "table.docx"))

    tables = [(p.page, b) for p in ir.pages for b in p.blocks if b.kind == "table"]
    assert len(tables) == 1
    page, block = tables[0]
    assert page == 1, f"table attributed to page {page}, but it sits on page 1"
    assert "Upstream" in block.text
    # ...and in reading order: after the heading it follows, not appended last.
    page_one = [b.text for b in ir.pages[0].blocks]
    assert page_one.index("Revenue Review") < page_one.index(block.text)


def test_a_slide_is_not_called_a_pdf_page(tmp_path):
    """A region title is user-facing text that ends up in a citation."""
    from src.unstructured.graph.axis1_structural import _PAGE_UNIT, _region_title

    assert _region_title("table", "", 3, 1, unit=_PAGE_UNIT["pptx"]) == "Table 1 (slide 3)"
    assert _region_title("table", "", 2, 1, unit=_PAGE_UNIT["xlsx"]) == "Table 1 (sheet 2)"
    assert _region_title("figure", "", 7, 2) == "Figure 2 (PDF page 7)"
